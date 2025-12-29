import os
import json
import re
from extract_pdf import extract_pdf_text
from clean_text import clean_text
from chunk_text import chunk_text
from pptx import Presentation
import nltk

nltk.download("punkt", quiet=True)

# ==============================
# UNIT MAPS
# ==============================

STALLINGS_UNIT_MAP = {
    1: "Unit 1", 2: "Unit 1",
    3: "Unit 2", 4: "Unit 2",
    5: "Unit 3", 6: "Unit 3",
    7: "Unit 4", 8: "Unit 4",
    9: "Unit 5"
}

KUROSE_UNIT_MAP = {
    1: "Unit 1", 2: "Unit 1",
    3: "Unit 2", 4: "Unit 2",
    5: "Unit 3",
    6: "Unit 4", 7: "Unit 4",
    8: "Unit 5"
}

# ==============================
# CONFIG
# ==============================

RAW_DIR = "data/raw"
OUTPUT_DIR = "data/processed/chunks"
os.makedirs(OUTPUT_DIR, exist_ok=True)

SKIP_KEYWORDS = [
    "edition", "copyright", "pearson",
    "table of contents", "preface",
    "associated companies", "global edition"
]

# ==============================
# HELPERS
# ==============================

def infer_chapter_from_section(text):
    match = re.search(r"\b(\d+)\.\d+\b", text)
    return int(match.group(1)) if match else None


def extract_chapter_number(text):
    # Try multiple patterns for chapter detection
    patterns = [
        r"\bchapter\s+(\d+)\b",
        r"^(\d+)\.\d+",  # Section numbers at start of line
        r"chapter\s*(\d+)",
        r"unit\s+(\d+)",
        r"^(\d+)\s+[A-Z]",  # Number followed by capital letter
        r"\n(\d+)\.\d+\s+[A-Z]",  # Section in new line
    ]
    
    for pattern in patterns:
        match = re.search(pattern, text, re.IGNORECASE | re.MULTILINE)
        if match:
            chapter = int(match.group(1))
            if 1 <= chapter <= 10:  # Valid chapter range
                return chapter
    return None


def get_unit_from_chapter(chapter, book_type):
    if book_type == "stallings":
        return STALLINGS_UNIT_MAP.get(chapter, "Unknown")
    if book_type == "kurose":
        return KUROSE_UNIT_MAP.get(chapter, "Unknown")
    return "Unknown"


def get_unit_from_filename_or_content(filename, content=""):
    """Map files to units based on filename or content keywords"""
    filename_lower = filename.lower()
    content_lower = content.lower()
    
    # Direct unit mention in filename
    match = re.search(r"unit\s*(\d+)", filename_lower)
    if match:
        return f"Unit {match.group(1)}"
    
    # Content-based mapping
    unit_keywords = {
        "Unit 1": ["data communication", "networking", "osi", "tcp/ip", "protocol", "introduction"],
        "Unit 2": ["physical", "data link", "ethernet", "csma", "token", "wireless", "bluetooth"],
        "Unit 3": ["network layer", "routing", "ip", "ipv4", "ipv6", "icmp", "dhcp", "nat"],
        "Unit 4": ["transport", "tcp", "udp", "application", "http", "ftp", "email", "dns"],
        "Unit 5": ["network management", "snmp", "monitoring", "wireshark", "sdn"]
    }
    
    # Check content for keywords
    for unit, keywords in unit_keywords.items():
        if any(keyword in content_lower for keyword in keywords):
            return unit
    
    # Filename-based fallback mapping
    if any(word in filename_lower for word in ["transport", "tcp", "udp"]):
        return "Unit 4"
    elif any(word in filename_lower for word in ["network", "layer"]):
        return "Unit 3"
    elif any(word in filename_lower for word in ["data", "link", "ethernet"]):
        return "Unit 2"
    elif any(word in filename_lower for word in ["email", "snmp"]):
        return "Unit 4" if "email" in filename_lower else "Unit 5"
    
    return "Unit 1"  # Default fallback


# ==============================
# PDF PROCESSOR (TEXTBOOKS)
# ==============================

def process_pdf_file(file_path, book_type):
    pages = extract_pdf_text(file_path)
    all_chunks = []
    
    current_unit = "Unit 1"  # Start with Unit 1 instead of Unknown
    current_chapter = None
    
    for page in pages:
        raw = page["text"].lower()
        if any(k in raw for k in SKIP_KEYWORDS):
            continue
            
        # Try to detect new chapter
        chapter = extract_chapter_number(page["text"])
        
        if not chapter:
            chapter = infer_chapter_from_section(page["text"])
            
        # Update unit if new chapter found
        if chapter and chapter != current_chapter:
            current_chapter = chapter
            new_unit = get_unit_from_chapter(chapter, book_type)
            if new_unit != "Unknown":
                current_unit = new_unit
        
        cleaned = clean_text(page["text"])
        if not cleaned.strip():  # Skip empty pages
            continue
            
        chunks = chunk_text(cleaned)
        
        for chunk in chunks:
            all_chunks.append({
                "unit": current_unit,
                "topic": f"Chapter {current_chapter}" if current_chapter else "General",
                "source": book_type,
                "page": page["page"],
                "text": chunk
            })
    
    return all_chunks


# ==============================
# NOTES
# ==============================

def process_notes_folder(notes_folder):
    all_chunks = []

    for pdf in os.listdir(notes_folder):
        if not pdf.lower().endswith(".pdf"):
            continue

        path = os.path.join(notes_folder, pdf)
        pages = extract_pdf_text(path)
        
        # Get sample content for unit detection
        sample_content = " ".join([p["text"][:500] for p in pages[:3]])
        unit = get_unit_from_filename_or_content(pdf, sample_content)

        for page in pages:
            cleaned = clean_text(page["text"])
            if not cleaned.strip():
                continue
                
            chunks = chunk_text(cleaned)

            for chunk in chunks:
                all_chunks.append({
                    "unit": unit,
                    "topic": "Notes",
                    "source": "notes",
                    "page": page["page"],
                    "text": chunk
                })

    return all_chunks


# ==============================
# PPTs
# ==============================

def process_ppts_folder(ppt_folder):
    all_chunks = []

    for file in os.listdir(ppt_folder):
        path = os.path.join(ppt_folder, file)
        
        # PDF slides
        if file.lower().endswith(".pdf"):
            pages = extract_pdf_text(path)
            sample_content = " ".join([p["text"][:500] for p in pages[:3]])
            unit = get_unit_from_filename_or_content(file, sample_content)
            
            for page in pages:
                cleaned = clean_text(page["text"])
                if not cleaned.strip():
                    continue
                    
                chunks = chunk_text(cleaned)

                for chunk in chunks:
                    all_chunks.append({
                        "unit": unit,
                        "topic": "Slide",
                        "source": "ppt",
                        "page": page["page"],
                        "text": chunk
                    })

        # PPTX
        elif file.lower().endswith(".pptx"):
            prs = Presentation(path)
            
            # Get sample content for unit detection
            sample_slides = []
            for slide in list(prs.slides)[:3]:
                text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
                sample_slides.append(" ".join(text))
            
            sample_content = " ".join(sample_slides)
            unit = get_unit_from_filename_or_content(file, sample_content)
            
            for i, slide in enumerate(prs.slides):
                text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)

                slide_text = "\n".join(text).strip()
                if slide_text:
                    all_chunks.append({
                        "unit": unit,
                        "topic": "Slide",
                        "source": "ppt",
                        "page": i + 1,
                        "text": slide_text
                    })

    return all_chunks


# ==============================
# SYLLABUS
# ==============================

def process_syllabus_file(file_path, units):
    pages = extract_pdf_text(file_path)
    syllabus_json = []

    pages_per_unit = max(1, len(pages) // len(units))

    for idx, page in enumerate(pages):
        unit = units[min(idx // pages_per_unit, len(units) - 1)]
        text = clean_text(page["text"])
        topics = [t.strip() for t in text.split("\n") if t.strip()]

        syllabus_json.append({
            "unit": unit,
            "topics": topics
        })

    out = os.path.join(OUTPUT_DIR, "syllabus.json")
    with open(out, "w", encoding="utf-8") as f:
        json.dump(syllabus_json, f, indent=2, ensure_ascii=False)

    print(f"[OK] Syllabus processed -> {out}")


# ==============================
# MAIN
# ==============================

# Create textbooks subdirectory
os.makedirs(os.path.join(OUTPUT_DIR, "textbooks"), exist_ok=True)

all_chunks = []
kurose_chunks = []
stallings_chunks = []

# 📘 TEXTBOOKS
textbooks_dir = os.path.join(RAW_DIR, "textbooks")
for book in os.listdir(textbooks_dir):
    book_type = "stallings" if "stallings" in book.lower() else "kurose"
    book_path = os.path.join(textbooks_dir, book)

    for pdf in os.listdir(book_path):
        if pdf.lower().endswith(".pdf"):
            chunks = process_pdf_file(os.path.join(book_path, pdf), book_type)
            if book_type == "kurose":
                kurose_chunks.extend(chunks)
            else:
                stallings_chunks.extend(chunks)
            all_chunks.extend(chunks)

# Save textbooks separately
with open(os.path.join(OUTPUT_DIR, "textbooks", "kurose.json"), "w", encoding="utf-8") as f:
    json.dump(kurose_chunks, f, indent=2, ensure_ascii=False)

with open(os.path.join(OUTPUT_DIR, "textbooks", "stallings.json"), "w", encoding="utf-8") as f:
    json.dump(stallings_chunks, f, indent=2, ensure_ascii=False)

print(f"[OK] Textbooks processed: Kurose({len(kurose_chunks)}), Stallings({len(stallings_chunks)})")

# NOTES
notes_chunks = process_notes_folder(os.path.join(RAW_DIR, "notes"))
all_chunks.extend(notes_chunks)

with open(os.path.join(OUTPUT_DIR, "notes.json"), "w", encoding="utf-8") as f:
    json.dump(notes_chunks, f, indent=2, ensure_ascii=False)

print(f"[OK] Notes processed: {len(notes_chunks)} chunks")

# PPTs
ppts_chunks = process_ppts_folder(os.path.join(RAW_DIR, "ppts"))
all_chunks.extend(ppts_chunks)

with open(os.path.join(OUTPUT_DIR, "ppts.json"), "w", encoding="utf-8") as f:
    json.dump(ppts_chunks, f, indent=2, ensure_ascii=False)

print(f"[OK] PPTs processed: {len(ppts_chunks)} chunks")

# SAVE ALL
with open(os.path.join(OUTPUT_DIR, "all_chunks.json"), "w", encoding="utf-8") as f:
    json.dump(all_chunks, f, indent=2, ensure_ascii=False)

print(f"[OK] All chunks saved: {len(all_chunks)} total chunks")

# 📑 SYLLABUS
syllabus_file = os.path.join(RAW_DIR, "syllabus", "syllabus.pdf")
process_syllabus_file(syllabus_file, ["Unit 1", "Unit 2", "Unit 3", "Unit 4", "Unit 5"])
